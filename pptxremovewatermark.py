import argparse
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def remove_hyperlinked_pictures_from_shapes(shapes_collection):
    """
    Iterates through a collection of shapes and removes pictures that have hyperlinks.
    """
    shapes_to_remove = []
    for i, shape in enumerate(shapes_collection):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            if shape.click_action and shape.click_action.hyperlink and shape.click_action.hyperlink.address:
                print(f"Found picture with hyperlink: {shape.click_action.hyperlink.address} on shape '{shape.name if shape.name else 'Unnamed Picture'}'")
                shapes_to_remove.append(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # Recursively check shapes within a group
            # Note: Modifying group shapes directly can be complex.
            # For simplicity, this example might not fully handle deeply nested groups
            # or might remove the entire group if a linked picture is found within.
            # A more robust solution would involve iterating group_items.
            # However, direct removal from group_items is not straightforward with python-pptx.
            # A common approach is to ungroup, remove, and regroup, which is complex.
            # For now, let's check if any sub-shape has a hyperlink.
            # This part is tricky as python-pptx doesn't allow easy removal of shapes *within* a group.
            # We will mark the group for removal if any of its picture items have a link.
            # This is a simplification.
            for sub_shape in shape.shapes: # shape.shapes for group shapes
                 if sub_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    if sub_shape.click_action and sub_shape.click_action.hyperlink and sub_shape.click_action.hyperlink.address:
                        print(f"Found picture with hyperlink within group '{shape.name}': {sub_shape.click_action.hyperlink.address}")
                        # Deciding to remove the whole group if a linked picture is inside.
                        # This might be too aggressive depending on the use case.
                        if shape not in shapes_to_remove:
                            shapes_to_remove.append(shape)
                        break # Found one, mark group and move to next main shape

    if shapes_to_remove:
        print(f"Attempting to remove {len(shapes_to_remove)} shapes...")
        for shape_to_remove in shapes_to_remove:
            try:
                element = shape_to_remove.element
                element.getparent().remove(element)
                print(f"Successfully removed shape.")
            except Exception as e:
                print(f"Could not remove shape: {e}")
    else:
        print("No hyperlinked pictures found in this collection.")


def process_presentation(input_path):
    """
    Opens a presentation, removes hyperlinked pictures from slide masters and layouts,
    and prompts the user to save.
    """
    try:
        prs = Presentation(input_path)
    except Exception as e:
        print(f"Error opening presentation '{input_path}': {e}")
        return

    print(f"Processing presentation: {input_path}")

    # Process Slide Masters
    if prs.slide_masters:
        print(f"\n--- Processing {len(prs.slide_masters)} Slide Master(s) ---")
        for i, slide_master in enumerate(prs.slide_masters):
            print(f"\nProcessing Slide Master {i+1}...")
            remove_hyperlinked_pictures_from_shapes(slide_master.shapes)

            # Process Slide Layouts within each master
            if slide_master.slide_layouts:
                print(f"\n--- Processing {len(slide_master.slide_layouts)} Slide Layout(s) for Master {i+1} ---")
                for j, slide_layout in enumerate(slide_master.slide_layouts):
                    print(f"\nProcessing Slide Layout {j+1} (from Master {i+1})...")
                    remove_hyperlinked_pictures_from_shapes(slide_layout.shapes)
            else:
                print(f"No slide layouts found for Master {i+1}.")
    else:
        print("No slide masters found in the presentation.")

    # Save the presentation
    while True:
        choice = input("Do you want to overwrite the original file (O) or save as a new file (S)? [O/S]: ").strip().upper()
        if choice == 'O':
            output_path = input_path
            break
        elif choice == 'S':
            base, ext = os.path.splitext(input_path)
            default_new_path = f"{base}_modified{ext}"
            new_path_input = input(f"Enter new file name (default: {default_new_path}): ").strip()
            output_path = new_path_input if new_path_input else default_new_path
            if os.path.abspath(output_path) == os.path.abspath(input_path):
                print("Error: New file name cannot be the same as the original if not overwriting. Please choose a different name.")
                continue
            break
        else:
            print("Invalid choice. Please enter 'O' or 'S'.")

    try:
        prs.save(output_path)
        print(f"Presentation saved to: {output_path}")
    except Exception as e:
        print(f"Error saving presentation to '{output_path}': {e}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Remove pictures with hyperlinks from PowerPoint slide masters and layouts.")
    parser.add_argument("input_file", help="Path to the input PowerPoint file (.pptx)")

    args = parser.parse_args()

    if not os.path.exists(args.input_file):
        print(f"Error: Input file '{args.input_file}' not found.")
    elif not args.input_file.lower().endswith(".pptx"):
        print(f"Error: Input file '{args.input_file}' is not a .pptx file.")
    else:
        process_presentation(args.input_file)
